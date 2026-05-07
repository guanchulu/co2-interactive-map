from __future__ import annotations

import re
import csv
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from export_joule_submission_docx import markdown_to_docx


ROOT = Path(__file__).resolve().parents[1]
SUBMISSION = ROOT / "docs" / "joule_submission"
INTERACTIVE = ROOT / "docs" / "interactive_map"
SI_MD = SUBMISSION / "supplemental_information.md"
SI_DOCX = SUBMISSION / "co2_supplemental_information.docx"
PPTX_OUT = SUBMISSION / "co2_all_figures_review_deck.pptx"
PPT_IMAGE_DIR = SUBMISSION / "figures_ppt"
FIGURE_MANIFEST = SUBMISSION / "co2_all_figures_review_manifest.csv"


def natural_key(path: Path) -> tuple[str, int, str]:
    match = re.search(r"figure(\d+)", path.stem, flags=re.IGNORECASE)
    number = int(match.group(1)) if match else 999
    return (str(path.parent), number, path.name)


def figure_inventory() -> list[tuple[str, Path]]:
    groups = [
        ("Main storyline", SUBMISSION / "figures_storyline", ["*.svg"]),
        ("Extended composite", SUBMISSION / "figures_composite", ["*.svg"]),
        ("Legacy diagnostic", SUBMISSION / "figures", ["*.svg", "*.png", "*.jpg", "*.jpeg"]),
        ("Interactive map preview", INTERACTIVE, ["preview*.png"]),
    ]
    items: list[tuple[str, Path]] = []
    seen: set[Path] = set()
    for group, folder, patterns in groups:
        if not folder.exists():
            continue
        paths: list[Path] = []
        for pattern in patterns:
            paths.extend(folder.glob(pattern))
        for path in sorted(paths, key=natural_key):
            resolved = path.resolve()
            if resolved in seen:
                continue
            seen.add(resolved)
            items.append((group, path))
    return items


def convert_svg_to_png(svg: Path, png: Path) -> Path:
    import cairosvg

    png.parent.mkdir(parents=True, exist_ok=True)
    cairosvg.svg2png(url=str(svg), write_to=str(png), output_width=2600)
    return png


def raster_for_ppt(path: Path) -> Path:
    suffix = path.suffix.lower()
    if suffix == ".svg":
        return convert_svg_to_png(path, PPT_IMAGE_DIR / f"{path.stem}.png")
    return path


def add_title(slide, text: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.35), Inches(0.2), Inches(12.65), Inches(0.55))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(20)
    paragraph.font.bold = True
    paragraph.alignment = PP_ALIGN.LEFT
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.35), Inches(0.72), Inches(12.65), Inches(0.32))
        p = sub.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(8)
        p.alignment = PP_ALIGN.LEFT


def add_fit_picture(slide, image_path: Path) -> None:
    with Image.open(image_path) as img:
        width_px, height_px = img.size
    max_left = Inches(0.35)
    max_top = Inches(1.1)
    max_width = Inches(12.65)
    max_height = Inches(6.0)
    image_ratio = width_px / height_px
    box_ratio = max_width / max_height
    if image_ratio >= box_ratio:
        width = max_width
        height = max_width / image_ratio
    else:
        height = max_height
        width = max_height * image_ratio
    left = max_left + (max_width - width) / 2
    top = max_top + (max_height - height) / 2
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def add_footer(slide, index: int, total: int, path: Path) -> None:
    footer = slide.shapes.add_textbox(Inches(0.35), Inches(7.15), Inches(12.65), Inches(0.22))
    p = footer.text_frame.paragraphs[0]
    p.text = f"{index}/{total}  {path.relative_to(ROOT)}"
    p.font.size = Pt(7)
    p.alignment = PP_ALIGN.RIGHT


def build_ppt() -> Path:
    items = figure_inventory()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    title_slide = prs.slides.add_slide(blank)
    add_title(title_slide, "CO2 allocation manuscript figure review deck")
    intro = title_slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.8), Inches(4.8))
    tf = intro.text_frame
    tf.text = f"{len(items)} figures are included. Each subsequent slide shows one figure at maximum readable size."
    tf.paragraphs[0].font.size = Pt(22)
    for group in sorted({group for group, _ in items}):
        p = tf.add_paragraph()
        p.text = f"{group}: {sum(1 for g, _ in items if g == group)}"
        p.font.size = Pt(16)

    total = len(items)
    for idx, (group, figure_path) in enumerate(items, start=1):
        slide = prs.slides.add_slide(blank)
        add_title(slide, f"{group}: {figure_path.stem}", str(figure_path.relative_to(ROOT)))
        raster = raster_for_ppt(figure_path)
        add_fit_picture(slide, raster)
        add_footer(slide, idx, total, figure_path)

    PPTX_OUT.parent.mkdir(parents=True, exist_ok=True)
    candidates = [PPTX_OUT] + [
        PPTX_OUT.with_name(f"{PPTX_OUT.stem}_v{i}{PPTX_OUT.suffix}") for i in range(2, 20)
    ]
    for candidate in candidates:
        try:
            prs.save(candidate)
            return candidate
        except PermissionError:
            continue
    raise PermissionError(f"Could not write any review deck candidate under {PPTX_OUT.parent}")


def write_manifest() -> Path:
    rows = []
    for idx, (group, figure_path) in enumerate(figure_inventory(), start=1):
        rows.append(
            {
                "slide_number": idx + 1,
                "figure_index": idx,
                "group": group,
                "name": figure_path.stem,
                "source_path": str(figure_path.relative_to(ROOT)),
            }
        )
    with FIGURE_MANIFEST.open("w", newline="", encoding="utf-8") as handle:
        writer = csv.DictWriter(handle, fieldnames=["slide_number", "figure_index", "group", "name", "source_path"])
        writer.writeheader()
        writer.writerows(rows)
    return FIGURE_MANIFEST


def main() -> None:
    markdown_to_docx(SI_MD, SI_DOCX)
    pptx = build_ppt()
    manifest = write_manifest()
    print(SI_DOCX)
    print(pptx)
    print(manifest)
    print(f"figures={len(figure_inventory())}")


if __name__ == "__main__":
    main()
